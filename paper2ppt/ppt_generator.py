"""
PowerPoint 生成模块 — 基于 python-pptx，支持自定义模板母版。

文字颜色、字体等样式继承自模板，仅按需调整字号与段落间距。
"""

from __future__ import annotations

import zipfile
from io import BytesIO
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt

from paper2ppt.slide_balance import estimate_content_font_size, estimate_title_font_size
from paper2ppt.table_extract import rows_look_valid
from paper2ppt.text_format import parse_inline_runs

LAYOUT_TITLE = 0
LAYOUT_CONTENT = 1
LAYOUT_SECTION = 2
LAYOUT_TITLE_ONLY = 5
LAYOUT_PICTURE = 8

# 幻灯片可用高度上限（英寸），预留页脚区域
SLIDE_HEIGHT_IN = 7.5
PICTURE_BOTTOM_MARGIN_IN = 0.55

TITLE_TYPES = frozenset({
    PP_PLACEHOLDER.TITLE,
    PP_PLACEHOLDER.CENTER_TITLE,
})
BODY_TYPES = frozenset({
    PP_PLACEHOLDER.BODY,
    PP_PLACEHOLDER.OBJECT,
})
SKIP_TYPES = frozenset({
    PP_PLACEHOLDER.PICTURE,
    PP_PLACEHOLDER.DATE,
    PP_PLACEHOLDER.FOOTER,
    PP_PLACEHOLDER.SLIDE_NUMBER,
    PP_PLACEHOLDER.HEADER,
})


def generate_ppt(
    slide_data: dict,
    output_path: str | Path,
    figures: dict[str, Path] | None = None,
    tables: dict[str, object] | None = None,
    template_path: str | Path | None = None,
) -> Path:
    output_path = Path(output_path)
    template_path = Path(template_path) if template_path else _default_template()

    if not template_path.exists():
        raise FileNotFoundError(f"PPT template not found: {template_path}")

    figure_bytes = _load_figure_bytes(figures or {})
    table_map = tables or {}

    prs = Presentation(str(template_path))
    _remove_existing_slides(prs)
    picture_layout_idx = _find_picture_layout(prs)

    inserted = 0
    table_count = 0
    for slide_info in slide_data.get("slides", []):
        slide_type = slide_info.get("type", "content")

        if slide_type == "title":
            _add_title_slide(prs, slide_info, slide_data)
        elif slide_type == "section":
            _add_section_slide(prs, slide_info)
        else:
            figure_id = slide_info.get("figure")
            image_data = figure_bytes.get(figure_id) if figure_id else None
            _add_content_slide(
                prs,
                slide_info,
                image_data=image_data,
                picture_layout_idx=picture_layout_idx,
                tables=table_map,
            )
            if image_data:
                inserted += 1
            if slide_info.get("table") or slide_info.get("table_data"):
                table_count += 1

    prs.save(str(output_path))
    print(f"      Embedded {inserted} figure(s), {table_count} table(s), {len(prs.slides)} slide(s) total")
    return output_path


def count_embedded_images(pptx_path: str | Path) -> int:
    with zipfile.ZipFile(pptx_path) as zf:
        return sum(1 for name in zf.namelist() if name.startswith("ppt/media/"))


def _default_template() -> Path:
    return Path(__file__).resolve().parent.parent / "template" / "index.pptx"


def _remove_existing_slides(prs: Presentation) -> None:
    slide_ids = list(prs.slides._sldIdLst)
    for sld_id in slide_ids:
        r_id = sld_id.rId
        prs.part.drop_rel(r_id)
        prs.slides._sldIdLst.remove(sld_id)


def _load_figure_bytes(figures: dict[str, Path]) -> dict[str, bytes]:
    loaded: dict[str, bytes] = {}
    for fid, path in figures.items():
        path = Path(path)
        if path.is_file():
            loaded[fid] = path.read_bytes()
    return loaded


def _find_picture_layout(prs: Presentation) -> int:
    for idx, layout in enumerate(prs.slide_layouts):
        for ph in layout.placeholders:
            if ph.placeholder_format.type == PP_PLACEHOLDER.PICTURE:
                return idx
    return min(LAYOUT_PICTURE, len(prs.slide_layouts) - 1)


def _find_title_placeholder(slide):
    """查找标题占位符（支持 TITLE 与 CENTER_TITLE）。"""
    if slide.shapes.title is not None:
        return slide.shapes.title
    for ph in slide.placeholders:
        if ph.placeholder_format.type in TITLE_TYPES:
            return ph
    return None


def _find_body_placeholder(slide):
    """
    查找正文占位符。

    模板「标题和内容」版式的正文类型为 OBJECT(7)，不是 BODY(2)；
    必须按类型匹配，绝不能回退到第一个非图片占位符（那会是标题）。
    """
    for ph in slide.placeholders:
        if ph.placeholder_format.type in BODY_TYPES:
            return ph
    return None


def _find_subtitle_placeholder(slide):
    for ph in slide.placeholders:
        if ph.placeholder_format.type == PP_PLACEHOLDER.SUBTITLE:
            return ph
    return None


def _configure_text_frame(text_frame) -> None:
    text_frame.word_wrap = True
    try:
        text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    except Exception:
        pass


def _set_font_size(shape, size: int) -> None:
    """仅调整字号，不覆盖模板颜色/字体。"""
    if shape is None:
        return
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(size)


def _add_title_slide(prs: Presentation, slide_info: dict, slide_data: dict):
    layout = prs.slide_layouts[min(LAYOUT_TITLE, len(prs.slide_layouts) - 1)]
    slide = prs.slides.add_slide(layout)

    title_text = slide_info.get("title") or slide_data.get("title", "Research Paper")
    title_ph = _find_title_placeholder(slide)
    if title_ph is not None:
        _set_formatted_text(
            title_ph.text_frame,
            title_text,
            estimate_title_font_size(title_text, default=36),
        )

    subtitle_ph = _find_subtitle_placeholder(slide)
    if subtitle_ph is not None:
        parts = []
        for value in (slide_info.get("subtitle"), slide_data.get("authors")):
            if value and value not in parts:
                parts.append(value)
        subtitle_ph.text = "\n".join(parts) if parts else ""
        _configure_text_frame(subtitle_ph.text_frame)
        _set_font_size(subtitle_ph, 18)

    return slide


def _add_section_slide(prs: Presentation, slide_info: dict):
    """生成章节过渡页，让听众明确讲述进度。"""
    layout = prs.slide_layouts[min(LAYOUT_SECTION, len(prs.slide_layouts) - 1)]
    slide = prs.slides.add_slide(layout)
    title = slide_info.get("title", "")
    title_ph = _find_title_placeholder(slide)
    if title_ph is not None:
        _set_formatted_text(title_ph.text_frame, title, estimate_title_font_size(title, 34))
    body = _find_body_placeholder(slide)
    if body is not None:
        body.text = slide_info.get("subtitle") or ""
    return slide


def _add_content_slide(
    prs: Presentation,
    slide_info: dict,
    *,
    image_data: bytes | None,
    picture_layout_idx: int,
    tables: dict[str, object] | None = None,
):
    has_figure = image_data is not None
    has_table = bool(slide_info.get("table") or slide_info.get("table_data"))
    if has_figure or has_table:
        return _add_visual_content_slide(
            prs,
            slide_info,
            image_data=image_data,
            tables=tables or {},
        )
    layout_idx = picture_layout_idx if has_figure else min(LAYOUT_CONTENT, len(prs.slide_layouts) - 1)
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])

    bullets = list(slide_info.get("bullets") or [])
    title_text = slide_info.get("title", "")

    title_ph = _find_title_placeholder(slide)
    if title_ph is not None:
        _set_formatted_text(title_ph.text_frame, title_text, estimate_title_font_size(title_text, default=28))

    body_ph = _find_body_placeholder(slide)
    font_size = estimate_content_font_size(bullets, has_figure or has_table)

    body_bottom = None
    if body_ph is not None:
        tf = body_ph.text_frame
        tf.clear()
        _configure_text_frame(tf)

        for i, bullet in enumerate(bullets):
            _add_bullet_lines(tf, bullet, font_size, is_first=(i == 0))

        if has_table and body_ph is not None:
            body_bottom = _insert_table_in_body(
                slide,
                slide_info,
                body_ph,
                tables or {},
                font_size=font_size,
            )

    if image_data:
        _insert_picture(slide, image_data)

    return slide


def _add_visual_content_slide(
    prs: Presentation,
    slide_info: dict,
    *,
    image_data: bytes | None,
    tables: dict[str, object],
):
    """图表页使用大画布：宽图/表格占满下半页，避免塞进狭窄右栏。"""
    layout = prs.slide_layouts[min(LAYOUT_TITLE_ONLY, len(prs.slide_layouts) - 1)]
    slide = prs.slides.add_slide(layout)
    title = slide_info.get("title", "")
    title_ph = _find_title_placeholder(slide)
    if title_ph is not None:
        _set_formatted_text(title_ph.text_frame, title, estimate_title_font_size(title, 28))

    bullets = list(slide_info.get("bullets") or [])
    has_table = bool(slide_info.get("table") or slide_info.get("table_data"))
    if has_table:
        body = slide.shapes.add_textbox(Inches(0.92), Inches(1.35), Inches(11.5), Inches(5.35))
        _fill_bullet_textbox(body, bullets, 17)
        _insert_table_in_body(slide, slide_info, body, tables, font_size=15)
        return slide

    img_w, img_h = _image_size_px(image_data or b"")
    aspect = img_w / img_h if img_h else 1.0
    if aspect >= 1.35:
        body = slide.shapes.add_textbox(Inches(0.92), Inches(1.33), Inches(11.5), Inches(1.55))
        _fill_bullet_textbox(body, bullets, 17)
        _add_picture_fitted(
            slide, image_data, Inches(0.92), Inches(3.0), Inches(11.5), Inches(3.65)
        )
    else:
        body = slide.shapes.add_textbox(Inches(0.92), Inches(1.45), Inches(6.05), Inches(5.15))
        _fill_bullet_textbox(body, bullets, 18)
        _add_picture_fitted(
            slide, image_data, Inches(7.25), Inches(1.45), Inches(5.15), Inches(5.15)
        )
    return slide


def _fill_bullet_textbox(shape, bullets: list[str], font_size: int) -> None:
    tf = shape.text_frame
    tf.clear()
    _configure_text_frame(tf)
    for index, bullet in enumerate(bullets):
        _add_bullet_lines(tf, bullet, font_size, is_first=(index == 0))


def _add_picture_fitted(slide, image_data: bytes | None, left, top, width, height) -> None:
    if not image_data:
        return
    img_w, img_h = _image_size_px(image_data)
    new_w, new_h = _fit_picture_emu(img_w, img_h, width, height)
    slide.shapes.add_picture(
        BytesIO(image_data),
        left + (width - new_w) // 2,
        top + (height - new_h) // 2,
        width=new_w,
        height=new_h,
    )


def _resolve_table_rows(slide_info: dict, tables: dict[str, object]) -> list[list[str]] | None:
    """优先 LLM 内联 table_data，其次 PDF 提取 rows。"""
    inline = slide_info.get("table_data")
    if isinstance(inline, dict):
        headers = inline.get("headers") or inline.get("header")
        rows = inline.get("rows") or []
        if headers:
            result = [list(headers)]
            result.extend([list(r) for r in rows if r])
            if len(result) >= 2:
                return result
        if rows and len(rows) >= 2:
            return [list(r) for r in rows]

    tid = slide_info.get("table")
    if tid and tid in tables:
        tbl = tables[tid]
        rows = getattr(tbl, "rows", None) or []
        if rows and rows_look_valid(rows):
            return rows
    return None


def _insert_table_in_body(
    slide,
    slide_info: dict,
    body_ph,
    tables: dict[str, object],
    *,
    font_size: int,
) -> int | None:
    """在正文区底部插入 PPT 原生表格；解析失败则嵌入 PDF 截取的表格图。"""
    rows = _resolve_table_rows(slide_info, tables)
    left = body_ph.left
    width = body_ph.width
    top = body_ph.top + int(body_ph.height * 0.52)
    max_height = body_ph.top + body_ph.height - Inches(0.08)
    avail_h = max_height - top
    if avail_h < Inches(0.5):
        top = body_ph.top + int(body_ph.height * 0.38)
        avail_h = max_height - top

    if rows and len(rows) >= 2:
        use_native = bool(slide_info.get("table_data")) or rows_look_valid(rows)
        if not use_native:
            rows = None

    if rows and len(rows) >= 2:
        row_count = len(rows)
        col_count = max(len(r) for r in rows)
        row_h = min(int(avail_h / row_count), Inches(0.32))
        table_h = row_h * row_count
        if table_h >= Inches(0.4):
            shape = slide.shapes.add_table(row_count, col_count, left, top, width, table_h)
            tbl = shape.table
            for ri, row in enumerate(rows):
                for ci in range(col_count):
                    cell_text = row[ci] if ci < len(row) else ""
                    cell = tbl.cell(ri, ci)
                    _set_formatted_text(
                        cell.text_frame,
                        str(cell_text),
                        max(font_size - 3, 8),
                        bold_first_line=(ri == 0),
                    )
            return top

    tid = slide_info.get("table")
    if tid and tid in tables:
        tbl = tables[tid]
        img_path = getattr(tbl, "image_path", None)
        if img_path and Path(img_path).is_file():
            data = Path(img_path).read_bytes()
            img_w, img_h = _image_size_px(data)
            pic_w = width
            pic_h = int(pic_w * img_h / img_w) if img_w else int(avail_h)
            pic_h = min(pic_h, avail_h)
            slide.shapes.add_picture(BytesIO(data), left, top, width=pic_w, height=pic_h)
            return top
    return None


def _set_formatted_text(
    text_frame,
    text: str,
    font_size: int,
    *,
    bold_first_line: bool = False,
) -> None:
    """向 text_frame 首段写入带 **粗体** / *斜体* 的文本。"""
    _configure_text_frame(text_frame)
    p = text_frame.paragraphs[0]
    p.text = ""
    for i, seg in enumerate(parse_inline_runs(text)):
        run = p.add_run()
        run.text = seg.text
        run.font.size = Pt(font_size)
        run.font.bold = seg.bold or (bold_first_line and i == 0)
        run.font.italic = seg.italic


def _set_formatted_paragraph(paragraph, text: str, font_size: int) -> None:
    """向 paragraph 写入 inline 格式。"""
    paragraph.text = ""
    for seg in parse_inline_runs(text):
        run = paragraph.add_run()
        run.text = seg.text
        run.font.size = Pt(max(font_size, 8))
        run.font.bold = seg.bold
        run.font.italic = seg.italic


def _add_bullet_lines(text_frame, bullet: str, font_size: int, *, is_first: bool) -> None:
    """添加一条 bullet，支持 \\n 多行与 **粗体** 等 inline 格式。"""
    lines = [line.strip() for line in str(bullet).split("\n") if line.strip()]
    if not lines:
        return

    for line_idx, line in enumerate(lines):
        if is_first and line_idx == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        p.level = 0 if line_idx == 0 else 1
        p.space_after = Pt(2 if line_idx else 4)
        size = max(font_size - (1 if line_idx else 0), 11)
        _set_formatted_paragraph(p, line, size)


def _image_size_px(image_data: bytes) -> tuple[int, int]:
    with Image.open(BytesIO(image_data)) as img:
        return img.size


def _fit_picture_emu(
    img_w: int,
    img_h: int,
    max_w: int,
    max_h: int,
) -> tuple[int, int]:
    """在限定框内等比缩放，返回 (width, height) EMU。"""
    if img_w <= 0 or img_h <= 0:
        return max_w, max_h

    width = max_w
    height = int(width * img_h / img_w)
    if height > max_h:
        height = max_h
        width = int(height * img_w / img_h)
    return width, height


def _insert_picture(slide, image_data: bytes) -> bool:
    img_w, img_h = _image_size_px(image_data)

    for ph in slide.placeholders:
        if ph.placeholder_format.type != PP_PLACEHOLDER.PICTURE:
            continue

        slot_left = ph.left
        slot_top = ph.top
        slot_width = ph.width
        slot_height = ph.height
        max_bottom = Inches(SLIDE_HEIGHT_IN - PICTURE_BOTTOM_MARGIN_IN)
        max_h = min(slot_height, max_bottom - slot_top)

        new_w, new_h = _fit_picture_emu(img_w, img_h, slot_width, max_h)
        new_left = slot_left + (slot_width - new_w) // 2
        new_top = slot_top + (slot_height - new_h) // 2

        slide.shapes.add_picture(
            BytesIO(image_data),
            new_left,
            new_top,
            width=new_w,
            height=new_h,
        )
        ph._sp.getparent().remove(ph._sp)
        return True

    slide.shapes.add_picture(
        BytesIO(image_data), Inches(7.0), Inches(1.5), width=Inches(5.5)
    )
    return True
